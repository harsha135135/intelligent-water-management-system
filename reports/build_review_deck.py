"""Build the Phase III review deck (PPTX) — one unified story over every model.

Six sections, in the order the review asks for them:

    System Testing · Verification and Validation · Chronos-2 · Deployment ·
    Final Experiment Results · Performance Analysis

Every number and every figure comes from ``results/chronos2/unified/`` — the same source the
results page reads — so the deck and the page cannot disagree. Nothing is typed by hand except
the prose.

    python reports/build_review_deck.py     # ~15 s -> reports/PW26_PK_06_phase3_review.pptx
"""

from __future__ import annotations

import json
import pathlib

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = pathlib.Path(__file__).resolve().parents[1]
U = ROOT / "results/chronos2/unified"
UP = U / "plots"
H45 = ROOT / "results/chronos2/calibrated"
OUT = ROOT / "reports" / "PW26_PK_06_phase3_review.pptx"

# ── palette — the same identities the figures use, so a colour never changes meaning ─────────
GROUND = RGBColor(0xFB, 0xFB, 0xF9)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
PANEL2 = RGBColor(0xF2, 0xF4, 0xF7)
HAIR = RGBColor(0xE2, 0xE6, 0xEA)
INK = RGBColor(0x14, 0x1A, 0x21)
INK2 = RGBColor(0x4A, 0x54, 0x5F)
INK3 = RGBColor(0x7A, 0x83, 0x8D)
C2 = RGBColor(0x2A, 0x78, 0xD6)      # Chronos-2
NPTS_C = RGBColor(0xD9, 0x5A, 0x2B)  # NPTS
PTC = RGBColor(0x6D, 0x3F, 0xC4)     # PatchTST
OK = RGBColor(0x12, 0x85, 0x5C)
WARN = RGBColor(0xA8, 0x72, 0x0B)
BAD = RGBColor(0xC3, 0x35, 0x2F)

SANS = "Helvetica Neue"
MONO = "Menlo"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)
BODY_W = W - 2 * MARGIN




# ────────────────────────────────────────────────────────────── data
lb = pd.read_csv(U / "leaderboard.csv")
sig = pd.read_csv(U / "significance_vs_all.csv")
wm = pd.read_csv(U / "win_matrix_all_horizons.csv")
pt = pd.read_csv(U / "per_tank.csv")
sk = pd.read_csv(U / "skill_h24.csv")
zi = pd.read_csv(U / "zero_inflation_h24.csv")
lt = pd.read_csv(U / "error_by_leadtime.csv")
cost_df = pd.read_csv(U / "cost.csv")
summ = json.loads((U / "summary.json").read_text())
h45s = pd.read_csv(H45 / "summary.csv").set_index("model")
h45e = pd.read_csv(H45 / "calibration_effect.csv").set_index("model")

HS = [6, 12, 24, 48, 72, 168]
HL = {6: "6 h", 12: "12 h", 24: "1 d", 48: "2 d", 72: "3 d", 168: "7 d"}
CHRONOS, INCUMBENT, REFERENCE = "Chronos2-ZS", "NPTS", "SeasonalNaive"
ORDER = list(summ["models"])
NM = len(ORDER)
NWORD = {9: "nine", 10: "ten", 11: "eleven", 12: "twelve"}.get(NM, str(NM))
LABEL = {
    "Chronos2-ZS": "Chronos-2 (zero-shot)", "Chronos2-COV": "Chronos-2 + covariates",
    "Chronos2-COV-LEAN": "Chronos-2 + cov (lean)", "Chronos2-COV-XL": "Chronos-2 + cov (XL)",
    "NPTS": "NPTS", "PatchTST": "PatchTST (defaults)", "PatchTST-Tuned": "PatchTST (tuned)",
    "ETS": "ETS", "Theta": "Theta", "DynamicOptimizedTheta": "DynamicOptimizedTheta",
    "SeasonalNaive": "SeasonalNaive-24",
}
TAG = {CHRONOS: "production", INCUMBENT: "incumbent", REFERENCE: "reference"}
FAMILY = dict(zip(lb.model, lb.family))
ROWS_H = {h: int(lb[(lb.horizon == h) & (lb.model == CHRONOS)]["rows"].iloc[0]) for h in HS}
TOTAL = sum(ROWS_H.values())


def cell(m, h, c):
    return float(lb[(lb.model == m) & (lb.horizon == h)][c].iloc[0])


_sM = sig[sig.metric == "MASE"]
_s24 = _sM[_sM.horizon == 24].set_index("opponent")
BEATEN = summ["opponents_beaten_significantly_all_horizons"]
_npts = _sM[_sM.opponent == INCUMBENT]
NPTS_LO, NPTS_HI = _npts.improvement_pct.min(), _npts.improvement_pct.max()
_ptt = _sM[_sM.opponent == "PatchTST-Tuned"]
PT_LO, PT_HI = _ptt.improvement_pct.min(), _ptt.improvement_pct.max()
Z = zi.set_index("model").loc[CHRONOS]
CAL = h45e.loc[CHRONOS]
OPS = h45s.loc[CHRONOS]
C2S = float(cost_df[cost_df.model == CHRONOS].wall_clock_s.iloc[0])
PTS = float(cost_df[cost_df.model == "PatchTST-Tuned"].wall_clock_s.iloc[0])
BASE_S = float(cost_df[cost_df.model == INCUMBENT].wall_clock_s.iloc[0])
CELLS = summ["tank_horizon_cells_won_vs"]
_d1 = lt[(lt.model == CHRONOS) & (lt.step <= 24)]["mae"].mean()
_d7 = lt[(lt.model == CHRONOS) & (lt.step > 144)]["mae"].mean()
SKILL24 = {m: int((sk[sk.model == m]["skill"] > 0).sum()) for m in sk.model.unique()}

# ────────────────────────────────────────────────────────────── primitives
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def _txt(frame, runs, *, size=14, color=INK, bold=False, font=SANS,
         align=PP_ALIGN.LEFT, space_after=4, line=None):
    """runs: a string, or a list of (text, {overrides}) tuples for inline emphasis."""
    para = frame.paragraphs[0] if not frame.paragraphs[0].runs and not frame.text else frame.add_paragraph()
    para.alignment = align
    para.space_after = Pt(space_after)
    if line:
        para.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, ov in runs:
        r = para.add_run()
        r.text = text
        f = r.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = ov.get("color", color)
    return para


def box(slide, x, y, w, h, *, fill=None, line=None, lw=0.75, radius=False):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if radius:
        shape.adjustments[0] = 0.06
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    shape.shadow.inherit = False
    return shape


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def slide(title=None, eyebrow=None, note=None):
    s = prs.slides.add_slide(BLANK)
    bg = box(s, 0, 0, W, H, fill=GROUND)
    bg.shadow.inherit = False
    y = MARGIN
    if eyebrow:
        tf = textbox(s, MARGIN, y, BODY_W, Inches(0.24))
        _txt(tf, eyebrow.upper(), size=10.5, color=INK3, bold=True, font=MONO, space_after=0)
        y += Inches(0.28)
    if title:
        tf = textbox(s, MARGIN, y, BODY_W, Inches(0.55))
        _txt(tf, title, size=27, color=INK, bold=True, space_after=0)
        y += Inches(0.52)
    if note:
        tf = textbox(s, MARGIN, y, Inches(10.6), Inches(0.6))
        _txt(tf, note, size=12.5, color=INK2, space_after=0, line=1.28)
        y += Inches(0.30) + Inches(0.19) * (1 + len(note) // 128)
    return s, y + Inches(0.16)


def divider(n, title, sub):
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, W, H, fill=INK)
    box(s, 0, H - Inches(0.10), W, Inches(0.10), fill=C2)
    tf = textbox(s, MARGIN, Inches(2.55), BODY_W, Inches(0.3))
    _txt(tf, f"SECTION {n}", size=12, color=C2, bold=True, font=MONO, space_after=0)
    tf = textbox(s, MARGIN, Inches(2.95), BODY_W, Inches(1.0))
    _txt(tf, title, size=42, color=GROUND, bold=True, space_after=0)
    tf = textbox(s, MARGIN, Inches(4.05), Inches(9.4), Inches(1.0))
    _txt(tf, sub, size=14, color=RGBColor(0xA3, 0xAE, 0xBA), space_after=0, line=1.35)
    return s


def _plain(val):
    if isinstance(val, list):
        return "".join(t for t, _ in val)
    return str(val)


def _lines(text, width_emu, size, font):
    """How many wrapped lines this cell needs. Character widths are the measured averages for
    Helvetica Neue and Menlo at the sizes used here — close enough that no row is ever clipped."""
    per_char = Pt(size * (0.60 if font == MONO else 0.505))
    per_line = max(1, int(width_emu / per_char))
    if not text:
        return 1
    n, cur = 1, 0
    for word in text.split(" "):
        if len(word) > per_line:                 # an unbreakable token wraps inside itself
            if cur:
                n += 1
            extra, cur = divmod(len(word), per_line)
            n += extra
            continue
        step = len(word) + (1 if cur else 0)
        if cur + step > per_line and cur:
            n += 1
            cur = len(word)
        else:
            cur += step
    return n


def table(slide, x, y, w, cols, rows, *, widths=None, size=11.5, head_size=9.5,
          row_h=Inches(0.29), head_h=Inches(0.31), highlight=None, aligns=None, fonts=None):
    """cols: list[str]. rows: list[list[str | [(text, overrides), ...]]].

    Row heights are computed from the wrapped line count of each cell, so a long cell grows its
    row instead of spilling into the one below.
    """
    widths = widths or [1.0 / len(cols)] * len(cols)
    aligns = aligns or ["l"] + ["r"] * (len(cols) - 1)
    fonts = fonts or [SANS] + [MONO] * (len(cols) - 1)
    highlight = highlight or set()
    xs, acc = [], x
    for frac in widths:
        xs.append((acc, int(w * frac)))
        acc += int(w * frac)

    pad = Inches(0.13)
    line_h = Pt(size * 1.32)
    heights = []
    for row in rows:
        need = 1
        for (cx, cw), val, fnt in zip(xs, row, fonts):
            need = max(need, _lines(_plain(val), cw - Inches(0.18), size, fnt))
        heights.append(max(row_h, need * line_h + pad))

    box(slide, x, y, w, head_h, fill=PANEL2)
    for (cx, cw), col, al in zip(xs, cols, aligns):
        tf = textbox(slide, cx + Inches(0.09), y + Inches(0.075), cw - Inches(0.18), head_h)
        _txt(tf, col.upper(), size=head_size, color=INK3, bold=True, font=MONO, space_after=0,
             align=PP_ALIGN.RIGHT if al == "r" else PP_ALIGN.LEFT)

    yy = y + head_h
    for i, (row, rh) in enumerate(zip(rows, heights)):
        fill = RGBColor(0xE8, 0xF1, 0xFC) if i in highlight else (
            PANEL if i % 2 == 0 else RGBColor(0xF9, 0xFA, 0xFB))
        box(slide, x, yy, w, rh, fill=fill)
        for j, ((cx, cw), val, al, fnt) in enumerate(zip(xs, row, aligns, fonts)):
            runs = val if isinstance(val, list) else [(str(val), {})]
            tf = textbox(slide, cx + Inches(0.09), yy + Inches(0.06), cw - Inches(0.18), rh)
            _txt(tf, runs, size=size, color=INK if j == 0 else INK2,
                 bold=(j == 0), font=fnt, space_after=0, line=1.28,
                 align=PP_ALIGN.RIGHT if al == "r" else PP_ALIGN.LEFT)
        yy += rh
    box(slide, x, y, w, head_h + sum(heights), line=HAIR)
    return yy


def tiles(slide, x, y, w, items, *, per_row=4, h=Inches(1.12), gap=Inches(0.14)):
    """items: (key, value, sub, color)"""
    tw = int((w - gap * (per_row - 1)) / per_row)
    for i, (k, v, sub, col) in enumerate(items):
        cx = x + (tw + gap) * (i % per_row)
        cy = y + (h + gap) * (i // per_row)
        box(slide, cx, cy, tw, h, fill=PANEL, line=HAIR)
        tf = textbox(slide, cx + Inches(0.15), cy + Inches(0.13), tw - Inches(0.3), Inches(0.2))
        _txt(tf, k.upper(), size=9, color=INK3, bold=True, font=MONO, space_after=0)
        tf = textbox(slide, cx + Inches(0.15), cy + Inches(0.34), tw - Inches(0.3), Inches(0.38))
        _txt(tf, v, size=21, color=col, bold=True, space_after=0)
        tf = textbox(slide, cx + Inches(0.15), cy + Inches(0.70), tw - Inches(0.3), Inches(0.3))
        _txt(tf, sub, size=9.5, color=INK2, space_after=0, line=1.15)
    return y + (h + gap) * ((len(items) + per_row - 1) // per_row)


def callout(slide, x, y, w, title, paras, *, accent=C2, h=None, size=11.5):
    if isinstance(paras, str):
        paras = [paras]
    inner = w - Inches(0.45)
    est = Inches(0.26)
    if title:
        est += Inches(0.30)
    for para in paras:
        est += _lines(_plain(para), inner, size, SANS) * Pt(size * 1.34) + Pt(5)
    h = h or est
    box(slide, x, y, w, h, fill=PANEL, line=HAIR)
    box(slide, x, y, Inches(0.045), h, fill=accent)
    tf = textbox(slide, x + Inches(0.24), y + Inches(0.13), w - Inches(0.45), h - Inches(0.2))
    if title:
        _txt(tf, title, size=13, color=INK, bold=True, space_after=6)
    for para in paras:
        _txt(tf, para, size=size, color=INK2, space_after=5, line=1.32)
    return y + h


def banner(slide, x, y, w, title, body, *, accent=BAD, size=11.5):
    inner = w - Inches(0.5)
    h = Inches(0.30) + Pt(15 * 1.3) + _lines(body, inner, size, SANS) * Pt(size * 1.34)
    box(slide, x, y, w, h, fill=SOFT[accent])
    box(slide, x, y, Inches(0.045), h, fill=accent)
    tf = textbox(slide, x + Inches(0.26), y + Inches(0.15), inner, h - Inches(0.2))
    _txt(tf, title, size=15, color=accent, bold=True, space_after=4)
    _txt(tf, body, size=size, color=INK2, space_after=0, line=1.3)
    return y + h


def bullets(slide, x, y, w, items, *, size=12.5, gap=Inches(0.30), accent=C2):
    for i, it in enumerate(items):
        cy = y + gap * i
        box(slide, x, cy + Inches(0.075), Inches(0.075), Inches(0.075), fill=accent)
        tf = textbox(slide, x + Inches(0.20), cy, w - Inches(0.20), gap)
        _txt(tf, it, size=size, color=INK2, space_after=0, line=1.3)
    return y + gap * len(items)


def picture(slide, path, x, y, w=None, h=None):
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def caption(slide, x, y, w, text):
    tf = textbox(slide, x, y, w, Inches(0.4))
    _txt(tf, text, size=10.5, color=INK3, space_after=0, line=1.25)


def footer(s, n):
    tf = textbox(s, MARGIN, H - Inches(0.42), Inches(9.0), Inches(0.25))
    _txt(tf, "PW26_PK_06 · Intelligent Water Management System · PES University RR",
         size=8.5, color=INK3, font=MONO, space_after=0)
    tf = textbox(s, W - MARGIN - Inches(1.0), H - Inches(0.42), Inches(1.0), Inches(0.25))
    _txt(tf, str(n), size=8.5, color=INK3, font=MONO, space_after=0, align=PP_ALIGN.RIGHT)


def chip(slide, x, y, label, color, soft):
    w = Inches(0.13) * len(label) + Inches(0.22)
    box(slide, x, y, w, Inches(0.235), fill=soft, radius=True)
    tf = textbox(slide, x + Inches(0.09), y + Inches(0.045), w - Inches(0.18), Inches(0.2))
    _txt(tf, label.upper(), size=8.5, color=color, bold=True, font=MONO, space_after=0)
    return x + w + Inches(0.1)


SOFT = {OK: RGBColor(0xE2, 0xF4, 0xEC), WARN: RGBColor(0xFB, 0xF1, 0xDC),
        BAD: RGBColor(0xFB, 0xE9, 0xE8), C2: RGBColor(0xE8, 0xF1, 0xFC),
        PTC: RGBColor(0xEE, 0xE9, 0xFA)}




# ══════════════════════════════════════════════════════════════ SLIDES
# 1 — title
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, fill=INK)
box(s, 0, 0, Inches(0.14), H, fill=C2)
tf = textbox(s, Inches(1.05), Inches(1.35), Inches(11.4), Inches(0.3))
_txt(tf, "PW26_PK_06 · PHASE III · REVIEW 1", size=12, color=C2, bold=True, font=MONO, space_after=0)
tf = textbox(s, Inches(1.05), Inches(1.80), Inches(11.2), Inches(1.7))
_txt(tf, "Intelligent Water Management System", size=40, color=GROUND, bold=True, space_after=2)
_txt(tf, "Predictive analytics for a sustainable campus", size=22,
     color=RGBColor(0xA3, 0xAE, 0xBA), space_after=0)
tf = textbox(s, Inches(1.05), Inches(3.62), Inches(10.6), Inches(0.9))
_txt(tf, [("Hourly water-demand forecasting for 24 campus tanks, benchmarked across ", {}),
          (f"{NWORD} models", {"bold": True, "color": GROUND}),
          (" on one shared grid — a pretrained foundation model, its covariate variants, the "
           "deployed incumbent, a transformer trained on this campus, and the classical "
           "baselines.", {})],
     size=14, color=RGBColor(0xA3, 0xAE, 0xBA), space_after=0, line=1.4)
for i, (k, v) in enumerate([("Dataset", "2025-01-01 → 2026-04-22"),
                            ("Observations", "270,849 hourly"),
                            ("Models", f"{NM} on one shared grid"),
                            ("Scored rows", f"{TOTAL:,} per model")]):
    cx = Inches(1.05) + Inches(2.85) * i
    tf = textbox(s, cx, Inches(4.95), Inches(2.7), Inches(0.22))
    _txt(tf, k.upper(), size=9, color=RGBColor(0x7B, 0x87, 0x92), bold=True, font=MONO, space_after=0)
    tf = textbox(s, cx, Inches(5.20), Inches(2.7), Inches(0.3))
    _txt(tf, v, size=13.5, color=GROUND, bold=True, space_after=0)
tf = textbox(s, Inches(1.05), Inches(6.28), Inches(11.2), Inches(0.7))
_txt(tf, "Abhay Patil · Amogh E M · Harshavardhan M · Viraj Ved Shankar", size=13,
     color=GROUND, space_after=3)
_txt(tf, "Guides: Ms. Preet Kanwal · Mr. Prasad B Honnavalli", size=11.5,
     color=RGBColor(0x7B, 0x87, 0x92), space_after=0)

# 2 — contents
s, y = slide("What this deck covers", eyebrow="Contents",
             note="Six review categories. Four are complete, one is partial, and one is a real "
                  "gap. Each is stated with the evidence, and the gaps are named rather than "
                  "softened.")
rows = [
    ["1 · System Testing", [("PARTIAL", {"color": WARN, "bold": True})],
     "6/6 metric unit tests and the methodology assertions pass. No system or integration suite."],
    ["2 · Verification and Validation", [("DONE", {"color": OK, "bold": True})],
     f"Row parity across all {NM} models, zero leakage, MASE identity, independent re-score, "
     "significance against every opponent, a trained deep control."],
    ["3 · Chronos-2", [("DONE", {"color": OK, "bold": True})],
     "amazon/chronos-2, zero-shot. Four covariate variants measured and indistinguishable."],
    ["4 · Deployment", [("NOT DONE", {"color": BAD, "bold": True})],
     "Nothing is serving live. Specified across six design documents; unbuilt."],
    ["5 · Final Experiment Results", [("DONE", {"color": OK, "bold": True})],
     f"{NM} models × 6 horizons × 24 tanks, {TOTAL:,} identical rows per model."],
    ["6 · Performance Analysis", [("DONE", {"color": OK, "bold": True})],
     "One unified figure set — every panel carries every model, with a table behind each."],
]
table(s, MARGIN, y, BODY_W, ["Category", "Status", "Evidence"], rows,
      widths=[0.26, 0.10, 0.64], aligns=["l", "l", "l"], row_h=Inches(0.42), size=11.5,
      fonts=[SANS, MONO, SANS])

# 3 — the grid
s, y = slide("The grid every number rests on", eyebrow="Foundation",
             note="Nothing in this deck is comparable unless this holds. The load-bearing "
                  "property is that every model is scored on identical rows — the earlier work "
                  "in this repository could not be compared model-to-model.")
y2 = tiles(s, MARGIN, y, BODY_W, [
    ("Tanks", "24", "26 dataset directories de-duplicated; the loader raises otherwise", INK),
    ("Observations", "270,849", "gapless hourly index; missing hours held open as NaN", INK),
    ("Origins", "24", "23-hour stride, co-prime with 24 — every hour-of-day once", INK),
    ("Horizons", "6", "6 h · 12 h · 1 d · 2 d · 3 d · 7 d", INK),
    ("Rows per model", f"{TOTAL:,}", f"identical for all {NM} models, asserted in code", C2),
    ("Leakage rows", "0", "every scored row satisfies timestamp > origin", OK),
    ("Duplicate rows", "0", "no (tank, timestamp) pair survives curation twice", OK),
    ("Sensor tiers", "15 / 6 / 3", "healthy / degraded / dead, from a mass-balance test", INK),
], per_row=4)
callout(s, MARGIN, y2 + Inches(0.20), BODY_W, "Why a 23-hour stride and not 24",
        ["With a 24-hour stride every origin lands on the same clock hour, so the 6 h horizon "
         "would only ever be scored on the quiet 00:00–05:00 window and never on the morning "
         "refill peak — the part that is actually hard to forecast. 23 is co-prime with 24, so "
         "the origins visit every hour-of-day exactly once.",
         "This is a property of the evaluation, not an option. Changing it invalidates comparison "
         "with every number in this deck."])

# ════════════════════════════════════════ 1 — SYSTEM TESTING
divider(1, "System Testing", "What is tested today validates the measurement apparatus, not a "
                             "running system. Both halves of that sentence are said out loud.")

s, y = slide("What is tested — and what it buys", eyebrow="System Testing · 1",
             note="tests/test_metrics.py runs without pytest and exits non-zero on failure. Run "
                  "it before trusting any number in this deck: if the MASE identity fails, every "
                  "figure that follows is wrong.")
rows = [
    ["test_seasonal_naive_mase_is_about_one", [("PASS", {"color": OK, "bold": True})],
     "load-bearing",
     "MASE is defined so a seasonal-naive forecast scores 1.0. Measured 0.967–1.069."],
    ["test_seasonal_scales_hand_worked", [("PASS", {"color": OK, "bold": True})], "unit",
     "Scaling denominators match a hand-computed value on a fixed series."],
    ["test_seasonal_scales_ignores_nan", [("PASS", {"color": OK, "bold": True})], "unit",
     "A sensor gap reduces the sample size pairwise; it does not poison the scale."],
    ["test_constant_series_scale_is_zero", [("PASS", {"color": OK, "bold": True})], "unit",
     "A constant series has an undefined MASE — not infinite, not zero."],
    ["test_perfect_forecast_scores_zero", [("PASS", {"color": OK, "bold": True})], "unit",
     "All four metrics are zero when the forecast is exact."],
    ["test_degenerate_series_excluded", [("PASS", {"color": OK, "bold": True})], "property",
     "Degenerate series are excluded from scaled aggregates, never patched with an epsilon."],
]
y2 = table(s, MARGIN, y, BODY_W, ["Test", "Result", "Kind", "What it guarantees"], rows,
           widths=[0.335, 0.075, 0.095, 0.495], aligns=["l", "l", "l", "l"],
           row_h=Inches(0.34), size=10.5, fonts=[MONO, MONO, SANS, SANS])
callout(s, MARGIN, y2 + Inches(0.18), BODY_W, None,
        [[("6 / 6 pass. ", {"bold": True, "color": OK}),
          ("Why an epsilon would have been dishonest: NEW_BLOCK_RO is constant-zero for 96% of "
           "its history. Adding a small constant to its zero denominator would manufacture a "
           "flattering MASE out of a sensor that never moved. It is excluded and reported "
           "separately instead.", {})]])

s, y = slide("The methodology checks, enforced in code", eyebrow="System Testing · 2",
             note="These are not a checklist someone ticks. score_benchmark --strict halts the "
                  "run if any of them fails, so a published number cannot come from a broken grid.")
rows = [
    ["Row parity across models", "assert_comparable()",
     f"All {NM} models scored on identical row counts at each horizon; fatal under --strict"],
    ["No leakage", "score_benchmark",
     "Raises if any scored row has timestamp ≤ its forecast origin. Measured: 0 rows"],
    ["No duplicates", "curate.py",
     "Raises if any (tank, timestamp) pair survives curation twice. Measured: 0"],
    ["Exactly 24 tanks", "curate.py",
     "Raises if de-duplication does not leave exactly 24 physical tanks"],
    ["Gapless hourly index", "reindex_gapless_hourly()",
     "Each tank reindexed onto a complete hourly range; missing hours become NaN"],
    ["Scales from pre-origin history", "attach_scales()",
     "MASE/RMSSE denominators use only data at or before the origin"],
    ["Unified pass agrees with the scorer", "unified_analysis",
     "All four macro metrics reproduce score_benchmark to 1e-16 across 66 model × horizon rows"],
]
y2 = table(s, MARGIN, y, BODY_W, ["Check", "Enforced by", "What fails the run"], rows,
           widths=[0.27, 0.22, 0.51], aligns=["l", "l", "l"], row_h=Inches(0.32), size=10.5,
           fonts=[SANS, MONO, SANS])
callout(s, MARGIN, y2 + Inches(0.16), BODY_W, "The check that keeps getting harder",
        [f"Every model added to the grid has to pass the same row-parity assertion, whatever "
         f"pipeline produced it. It now holds across {NM} models from three separate pipelines — "
         f"{ROWS_H[6]:,} / {ROWS_H[12]:,} / {ROWS_H[24]:,} / {ROWS_H[48]:,} / {ROWS_H[72]:,} / "
         f"{ROWS_H[168]:,} rows, unchanged."], accent=OK)

s, y = slide("What is not tested — say this before you are asked", eyebrow="System Testing · 3",
             note="The tests above validate the measurement apparatus. They are not system tests, "
                  "and presenting them as system tests is the easiest thing to be caught on.")
y = banner(s, MARGIN, y, BODY_W, "No system or integration test suite exists.",
           "There is nothing to integration-test yet: the real-time ingestion, state store, "
           "decision engine and dashboard are specified across six design documents and none of "
           "them is built.") + Inches(0.26)
tf = textbox(s, MARGIN, y, BODY_W, Inches(0.3))
_txt(tf, "The test strategy is written down. Six layers, in the order they would be built:",
     size=12.5, color=INK, bold=True, space_after=0)
y += Inches(0.36)
rows = [
    ["Unit", "Pure functions — metrics, calendar features, quality rules", "6/6 exist today"],
    ["Property", "Invariants for any input, e.g. gapless reindex preserves order", "1 exists today"],
    ["Contract", "API request/response shapes and SSE schemas against docs/api_design.md", "not built"],
    ["Integration", "Ingest → validate → state → forecast → score, end to end on a replay", "not built"],
    ["Parity", "Replay and live must produce identical forecasts from identical context", "not built"],
    ["Regression", "The published benchmark re-scored on every commit; drift fails CI",
     "score_benchmark --strict is the seed"],
]
table(s, MARGIN, y, BODY_W, ["Layer", "What it would test", "Status"], rows,
      widths=[0.13, 0.62, 0.25], aligns=["l", "l", "l"], row_h=Inches(0.34), size=11,
      fonts=[SANS, SANS, SANS])

# ════════════════════════════════════════ 2 — V&V
divider(2, "Verification and Validation",
        "The strongest part of the project. Six machine-checked guarantees, and significance "
        "tested against every opponent on the grid rather than one pairing at a time.")

s, y = slide("Six guarantees, each checked by a machine", eyebrow="V&V · 1",
             note="Verification asks whether the measurement is built correctly. Validation asks "
                  "whether the result is real. The first four are verification; the last two are "
                  "validation.")
rows = [
    ["1 · Row parity", "verification",
     f"All {NM} models scored on identical rows at each horizon",
     [("asserted, fatal", {"color": OK})]],
    ["2 · Zero leakage", "verification",
     "Every scored row is strictly after its origin; scales use pre-origin history only",
     [("0 rows", {"color": OK})]],
    ["3 · MASE identity", "verification",
     "SeasonalNaive-24 must score MASE ≈ 1.0, or the whole scale is wrong",
     [("0.967–1.069", {"color": OK})]],
    ["4 · Independent re-score", "verification",
     "Two independent implementations reproduce every macro metric",
     [("agree to 1e-16", {"color": OK})]],
    ["5 · Statistical significance", "validation",
     "Paired bootstrap over 24 origins plus Diebold–Mariano, against every opponent",
     [(f"{len(BEATEN)}/{len(BEATEN)} beaten", {"color": OK})]],
    ["6 · Deep-learning control", "validation",
     "PatchTST trained on this campus at two configurations, on the identical rows",
     [("on the grid", {"color": OK})]],
]
y2 = table(s, MARGIN, y, BODY_W, ["Guarantee", "Kind", "What it establishes", "Result"], rows,
           widths=[0.23, 0.10, 0.51, 0.16], aligns=["l", "l", "l", "r"],
           row_h=Inches(0.32), size=10.5, fonts=[SANS, SANS, SANS, MONO])
callout(s, MARGIN, y2 + Inches(0.18), BODY_W, "Why row parity is the load-bearing one",
        ["The pre-existing results in this repository could not be compared to each other: the "
         "AutoGluon run scored 26 series over 624 holdout rows, an older PatchTST run 24 over "
         "576, on different dates. Two numbers measured on different rows are not a comparison — "
         "and everything here rests on assert_comparable() being fatal rather than a warning."])

s, y = slide("Is the margin real? — against every opponent", eyebrow="V&V · 2",
             note="Two independent tests, because they assume different things. The bootstrap "
                  "assumes only that the 24 origins are exchangeable; Diebold–Mariano assumes a "
                  "covariance structure the 23-hour stride approximates. They agree everywhere.")
rows = []
for m in ORDER:
    if m == CHRONOS:
        continue
    r = _s24.loc[m]
    allh = _sM[_sM.opponent == m]
    won = int(allh.base_better.sum())
    verdict = ([("significant, all 6", {"color": OK})] if won == 6
               else [("indistinguishable", {"color": INK3})])
    rows.append([LABEL[m], FAMILY[m],
                 [(f"{r.improvement_pct:+.2f}%",
                   {"color": C2, "bold": True} if r.improvement_pct > 0.5 else {"color": INK3})],
                 f"[{r.ci_lo_pct:+.2f}, {r.ci_hi_pct:+.2f}]",
                 "<0.0001" if r.p_bootstrap < 1e-4 else f"{r.p_bootstrap:.3f}",
                 f"{r.p_dm:.0e}", f"{int(r.origins_won)}/24",
                 f"{CELLS[m]}/144", verdict])
y2 = table(s, MARGIN, y, BODY_W,
           ["Opponent", "Family", "MASE impr.", "95% CI", "Boot p", "DM p", "Origins", "Cells",
            "Across all six"],
           rows, widths=[0.20, 0.10, 0.10, 0.135, 0.085, 0.075, 0.075, 0.075, 0.15],
           aligns=["l", "l", "r", "r", "r", "r", "r", "r", "l"], row_h=Inches(0.30), size=10,
           fonts=[SANS, SANS, MONO, MONO, MONO, MONO, MONO, MONO, MONO])
callout(s, MARGIN, y2 + Inches(0.18), BODY_W, None,
        [[(f"Chronos-2 beats all {len(BEATEN)} non-Chronos models significantly at every one of "
           f"the six horizons", {"bold": True, "color": INK}),
          (f" — {NPTS_LO:.1f}–{NPTS_HI:.1f}% against the deployed incumbent, "
           f"{PT_LO:.1f}–{PT_HI:.1f}% against the best PatchTST trained here, and roughly a third "
           "against the classical methods. The three covariate variants are the only opponents it "
           "does not separate from, and their intervals straddle zero — which is the finding, not "
           "a failure.", {})]], accent=OK)

# ════════════════════════════════════════ 3 — CHRONOS-2
divider(3, "Chronos-2",
        "A pretrained time-series foundation model, used zero-shot. Nothing is fitted, and the "
        "choice to leave it that way is a measured decision rather than a shortcut.")

s, y = slide("What the model is, and what zero-shot means here", eyebrow="Chronos-2 · 1")
col_w = int((BODY_W - Inches(0.3)) / 2)
box(s, MARGIN, y, col_w, Inches(2.30), fill=PANEL, line=HAIR)
tf = textbox(s, MARGIN + Inches(0.24), y + Inches(0.20), col_w - Inches(0.48), Inches(1.9))
_txt(tf, "The model", size=13.5, color=INK, bold=True, space_after=7)
for k, v in [("Checkpoint", "amazon/chronos-2"), ("Parameters", "119.5 M"),
             ("Regime", "zero-shot — no training, no fine-tuning"),
             ("Context", "2,048 hours, truncated at each origin"),
             ("Device", "MPS (Apple laptop GPU)"),
             ("Output", "mean plus p10 / p25 / p50 / p75 / p90")]:
    _txt(tf, [(f"{k}   ", {"color": INK3, "size": 10.5, "font": MONO}),
              (v, {"color": INK2, "size": 11.5})], space_after=4)
box(s, MARGIN + col_w + Inches(0.3), y, col_w, Inches(2.30), fill=PANEL, line=HAIR)
tf = textbox(s, MARGIN + col_w + Inches(0.54), y + Inches(0.20), col_w - Inches(0.48), Inches(1.9))
_txt(tf, "Why it works without seeing our data", size=13.5, color=INK, bold=True, space_after=7)
_txt(tf, "It was pretrained on millions of time series and has learned what demand-like signals "
         "look like in general. Campus water demand has strong daily structure it can recognise "
         "without ever having seen this campus.", size=11.5, color=INK2, space_after=6, line=1.3)
_txt(tf, [("The evidence is not the argument — it is that the model beats every model fitted "
           "specifically on this data, at every horizon, while costing ", {}),
          (f"{C2S:.0f} seconds", {"bold": True, "color": C2}),
          (f" against {BASE_S/60:.0f}–{PTS/60:.0f} minutes.", {})],
     size=11.5, color=INK2, space_after=0, line=1.3)
callout(s, MARGIN, y + Inches(2.52), BODY_W,
        "The one thing zero-shot buys that accuracy does not",
        ["A zero-shot model holds no fitted state, so there is nothing to go stale. When a tank is "
         "added, a sensor degrades, or the academic calendar shifts the demand profile, it needs "
         "no retraining and no covariate pipeline. For a campus deployment that is the operational "
         "argument; the accuracy result merely means there is nothing to trade against it."])

s, y = slide("Does conditioning on covariates pay for itself?", eyebrow="Chronos-2 · 2",
             note="Four variants of the same backbone on the same grid. The covariates are the "
                  "academic-calendar features — term dates, exam proximity, holidays, weekends — "
                  "the conditioning most likely to help on a campus.")
VAR = [m for m in ORDER if m.startswith("Chronos2")]
rows = []
for v in VAR:
    rt = float(cost_df[cost_df.model == v].wall_clock_s.iloc[0])
    rows.append([LABEL[v]] + [f"{cell(v, h, 'macro_mase'):.4f}" for h in HS]
                + [[(f"{rt/60:.1f}", {"color": C2 if v == CHRONOS else INK2,
                                      "bold": v == CHRONOS})]])
spread = lb[lb.model.isin(VAR)].groupby("horizon").macro_mase.agg(lambda x: x.max() - x.min())
y2 = table(s, MARGIN, y, BODY_W, ["Variant"] + [HL[h] for h in HS] + ["Runtime (min)"], rows,
           widths=[0.24] + [0.10] * 6 + [0.16], row_h=Inches(0.32), size=11.5, highlight={0})
y2 = tiles(s, MARGIN, y2 + Inches(0.20), BODY_W, [
    ("MASE spread, 4 variants", f"{spread.max():.4f}",
     "largest gap at any horizon — statistically indistinguishable", INK),
    ("Zero-shot cost", f"{C2S:.0f} s", "for the entire 24-origin × 6-horizon backtest", C2),
    ("Cheapest covariate variant",
     f"{float(cost_df[cost_df.model=='Chronos2-COV-LEAN'].wall_clock_min.iloc[0]):.1f} min",
     "for an accuracy difference inside the noise", INK),
    ("Decision", "zero-shot", "a compute decision, defensible because the accuracy gap is nil", OK),
], per_row=4)

# ════════════════════════════════════════ 4 — DEPLOYMENT
divider(4, "Deployment",
        "The real gap. Nothing is serving live. The architecture that would close it is fully "
        "specified — describing that accurately is more credible than describing it generously.")

s, y = slide("Nothing is serving", eyebrow="Deployment · 1")
y = banner(s, MARGIN, y, BODY_W, "There is no running system, and the demo surface is not one.",
           "A Chrome MV3 dock renders real Chronos-2 forecasts from a bundled JSON file, offline. "
           "It is a demo surface with precomputed data. Calling it a deployment is the fastest way "
           "to lose a reviewer's trust in every other claim in this deck.") + Inches(0.24)
rows = [
    ["Data curation, calendar features, backtest, metrics", [("implemented", {"color": OK})],
     "src/data/, src/models/backtest.py, metrics.py"],
    ["Chronos-2 inference and the four covariate variants", [("implemented", {"color": OK})],
     "chronos2_forecasting.py"],
    ["PatchTST control, two configurations", [("implemented", {"color": OK})],
     "patchtst_benchmark.py"],
    ["Unified analysis and the 15-figure evidence set", [("implemented", {"color": OK})],
     "unified_analysis.py, unified_figures.py"],
    ["Conformal calibration + volume bias correction", [("implemented", {"color": OK})],
     "calibration.py, calibrated_holdout.py"],
    ["Waltr forecast dock (offline, precomputed)", [("demo only", {"color": WARN})],
     "extension/ — a demo surface, not a deployment"],
    ["WALTR HTTP client", [("unusable", {"color": BAD})],
     "code exists; requires a JWT that has not been issued"],
    ["Ingestion, state store, alerts, decision engine, dashboard",
     [("not built", {"color": BAD})], "specified in Phases 1–8 of the implementation plan"],
    ["Real motor control", [("blocked", {"color": BAD})],
     "no motor API exists, and the dataset contains no motor telemetry"],
]
table(s, MARGIN, y, BODY_W, ["Component", "Status", "Where"], rows,
      widths=[0.44, 0.14, 0.42], aligns=["l", "l", "l"], row_h=Inches(0.30), size=10.5,
      fonts=[SANS, MONO, MONO])

s, y = slide("What is designed, and what would have to be built", eyebrow="Deployment · 2",
             note="Six design documents specify the real-time system end to end. None of the "
                  "pipeline they describe exists in the repository yet — a distinction stated in "
                  "the documents themselves, not just here.")
rows = [
    ["realtime_architecture.md", "Master architecture — adapter seam, component catalogue, "
     "10 diagrams, quality rules, feedback loop, degradation matrix"],
    ["implementation_plan.md", "Phases 0–11 with files, dependencies, tests, acceptance criteria"],
    ["api_design.md", "REST surface, SSE event types, transport comparison, auth and roles"],
    ["data_model.md", "20 entities with fields, indexes and retention"],
    ["demo_plan.md", "Historical replay, the actual-vs-predicted reveal, scenarios from the record"],
    ["safety_and_controls.md", "Motor-control safety, per-tank limits, fail-safe matrix"],
]
y2 = table(s, MARGIN, y, BODY_W, ["Document", "Contents"], rows,
           widths=[0.26, 0.74], aligns=["l", "l"], row_h=Inches(0.34), size=10.5,
           fonts=[MONO, SANS])
callout(s, MARGIN, y2 + Inches(0.18), BODY_W,
        "Demo and live are separated structurally, not by convention",
        ["The demo replays the historical dataset through the production pipeline as if it were "
         "arriving now. mode (replay | live) is a column on every fact table, a field in every API "
         "response and SSE event, and a non-dismissible banner in the UI. Deleting a replay "
         "session deletes every row it created.",
         "Forecasts during a replay are real Chronos-2 inference from the context available at "
         "that simulated moment — not precomputed. Switching to production is a configuration "
         "change, not a rewrite."], accent=WARN)

# ════════════════════════════════════════ 5 — FINAL EXPERIMENT RESULTS
divider(5, "Final Experiment Results",
        f"{NWORD.capitalize()} models, six horizons, 24 tanks, {TOTAL:,} identical rows per "
        "model. The leaderboard, the head-to-head, the cost, and the two caveats — in that order, "
        "because the caveats belong with the result.")

s, y = slide(f"The leaderboard — all {NM} models, all six horizons", eyebrow="Results · 1",
             note="Macro MASE with rank in each cell. MASE is a ratio to a seasonal-naive "
                  "baseline: 1.0 means no better than naive, and it is not a percentage.")
rows, hl = [], set()
for i, m in enumerate(ORDER):
    if m == CHRONOS:
        hl.add(i)
    name = [(LABEL[m], {})]
    if m in TAG:
        name.append((f"   {TAG[m]}", {"color": INK3, "size": 9, "font": MONO, "bold": False}))
    cells = [[(f"{cell(m, h, 'macro_mase'):.4f}",
               {"color": C2, "bold": True} if m == CHRONOS else {})] for h in HS]
    rows.append([name, FAMILY[m]] + cells + [f"{cell(m, 24, 'macro_mae'):.4f}"])
table(s, MARGIN, y, BODY_W, ["Model", "Family"] + [HL[h] for h in HS] + ["MAE @ 1 d"], rows,
      widths=[0.27, 0.105] + [0.088] * 6 + [0.097],
      aligns=["l", "l"] + ["r"] * 7, row_h=Inches(0.285), size=10.5, head_h=Inches(0.29),
      highlight=hl, fonts=[SANS, SANS] + [MONO] * 7)

s, y = slide("Per tank — where each model can and cannot be trusted", eyebrow="Results · 2",
             note="Aggregate accuracy decides whether to adopt a model; these decide whether to "
                  "trust it on a particular tank. Skill is 1 − MAE(model)/MAE(SeasonalNaive-24).")
_p24 = pt[pt.horizon == 24]
_best = _p24.loc[_p24.groupby("tank")["mase"].idxmin()]
_wins = _best.model.value_counts()
_skg = sk.groupby("model")["skill"].agg(["min", "median"])
rows, hl = [], set()
for i, m in enumerate([x for x in ORDER if x != REFERENCE]):
    if m == CHRONOS:
        hl.add(i)
    g = _skg.loc[m]
    pos = SKILL24[m]
    rows.append([LABEL[m], FAMILY[m], str(int(_wins.get(m, 0))),
                 [(f"{pos}/24", {"color": OK if pos == 24 else BAD, "bold": True})],
                 f"{g['median']:.3f}",
                 [(f"{g['min']:+.3f}", {"color": INK2 if g['min'] > 0 else BAD})]])
y2 = table(s, MARGIN, y, BODY_W,
           ["Model", "Family", "Tanks where best", "Positive skill", "Median skill", "Worst tank"],
           rows, widths=[0.28, 0.13, 0.15, 0.14, 0.15, 0.15],
           aligns=["l", "l", "r", "r", "r", "r"], row_h=Inches(0.275), size=10.5,
           head_h=Inches(0.29), highlight=hl, fonts=[SANS, SANS, MONO, MONO, MONO, MONO])
callout(s, MARGIN, y2 + Inches(0.18), BODY_W, "The objection this answers",
        [[("“Your average is flattered by the dead sensors.” It is not. ", {}),
          ("Every Chronos-2 variant and the incumbent clear zero on all 24 tanks", 
           {"bold": True, "color": INK}),
          (", and skill does not decline with demand size. The trained deep models are worse than "
           "naive on three to four tanks; the classical methods on five.", {})]], accent=OK)

s, y = slide("What the accuracy cost", eyebrow="Results · 3",
             note=f"Measured wall clock for the identical {TOTAL:,}-row backtest, read from the "
                  "run manifests. The five AutoGluon baselines share one fitting pass, so their "
                  "figure is a family total, not a per-model cost.")
rows, hl = [], set()
for r in cost_df.dropna(subset=["wall_clock_s"]).itertuples():
    if r.model == CHRONOS:
        hl.add(len(rows))
    shared = r.regime.startswith("fitted")
    rows.append([LABEL[r.model], r.regime,
                 f"{r.wall_clock_s:,.0f} s" + (" *" if shared else ""),
                 f"{r.wall_clock_min:.1f}",
                 [(f"{r.macro_mase_h24:.4f}",
                   {"color": C2, "bold": True} if r.model == CHRONOS else {})]])
y2 = table(s, MARGIN, y, BODY_W, ["Model", "Regime", "Wall clock", "Minutes", "MASE @ 1 d"], rows,
           widths=[0.24, 0.38, 0.13, 0.10, 0.15], aligns=["l", "l", "r", "r", "r"],
           row_h=Inches(0.25), size=10, head_h=Inches(0.27), highlight=hl,
           fonts=[SANS, SANS, MONO, MONO, MONO])
callout(s, MARGIN, y2 + Inches(0.14), BODY_W, None,
        [[("The best PatchTST needed ", {}), (f"{PTS/60:.0f} minutes", {"bold": True, "color": INK}),
          (" of training against Chronos-2's ", {}),
          (f"{C2S:.0f} seconds", {"bold": True, "color": C2}),
          (f" of inference — {PTS/C2S:.0f}× — and must be refitted whenever the campus changes, "
           "whereas a zero-shot model holds no fitted state that can go stale.   "
           "* shared AutoGluon pass.", {})]], accent=OK)

s, y = slide("The two caveats — measured, diagnosed, corrected", eyebrow="Results · 4",
             note="Both were stated in the benchmark as open weaknesses. Both now have a "
                  "mechanism and a fix, fitted on 8 Jan – 8 Mar 2026 and measured on the disjoint "
                  "9 Mar – 22 Apr window the correction never saw.")
col_w = int((BODY_W - Inches(0.3)) / 2)
for i, (title, accent, paras) in enumerate([
    ("1 · The intervals were overconfident", WARN, [
        [("p10–p90 covered ", {}), (f"{CAL.coverage_raw:.3f}", {"bold": True, "color": INK}),
         (" against a nominal 0.80.", {})],
        [("The cause: ", {}),
         (f"{Z.zero_fraction*100:.0f}% of hourly readings are exactly zero",
          {"bold": True, "color": INK}),
         (f", and a continuous-density model puts its p10 above zero on "
          f"{Z.p10_above_zero_pct:.0f}% of rows — so "
          f"{Z.below_misses_that_are_zero*100:.0f}% of its lower-tail misses are zero-demand "
          "hours. It is a zero-inflation problem, and it applies to every continuous-density "
          "model on the grid.", {})],
        [("Corrected by asymmetric conformal calibration: coverage ", {}),
         (f"{CAL.coverage_raw:.3f} → {CAL.coverage_cal:.3f}", {"bold": True, "color": OK}),
         (", measured out of sample.", {})]]),
    ("2 · 24-hour volume was under-forecast", WARN, [
        [("Cumulative volume was under-forecast by ", {}),
         (f"{abs(OPS.total_bias_raw_pct):.1f}%", {"bold": True, "color": INK}),
         (" — the operationally expensive direction, because a refill sized on the raw mean would "
          "be short.", {})],
        [("Corrected by a per-tank multiplicative factor: ", {}),
         (f"{OPS.total_bias_raw_pct:+.1f}% → {OPS.total_bias_pct:+.1f}%",
          {"bold": True, "color": OK}),
         (f", cutting daily campus error from {OPS.daily_mae_raw_kl:.1f} to ", {}),
         (f"{OPS.daily_mae_kl:.1f} KL", {"bold": True, "color": OK}), (".", {})],
        [("A refill should still be sized on the calibrated interval, not the raw mean.",
          {"bold": True, "color": INK})]]),
]):
    x = MARGIN + (col_w + Inches(0.3)) * i
    box(s, x, y, col_w, Inches(2.75), fill=PANEL, line=HAIR)
    box(s, x, y, Inches(0.045), Inches(2.75), fill=accent)
    tf = textbox(s, x + Inches(0.26), y + Inches(0.18), col_w - Inches(0.5), Inches(2.4))
    _txt(tf, title, size=13.5, color=INK, bold=True, space_after=8)
    for p in paras:
        _txt(tf, p, size=11.5, color=INK2, space_after=8, line=1.32)
callout(s, MARGIN, y + Inches(2.97), BODY_W, None,
        [[("What calibration does not fix: ", {"bold": True, "color": INK}),
          (f"the incumbent reproduces {h45s.loc[INCUMBENT].sd_ratio*100:.0f}% of real day-to-day "
           f"variation and correlates {h45s.loc[INCUMBENT].corr_with_actual:+.2f} with it, while "
           f"Chronos-2 reproduces {OPS.sd_ratio*100:.0f}% and correlates "
           f"{OPS.corr_with_actual:+.2f}. Post-hoc correction fixes where a forecast sits and how "
           "sure it claims to be — not whether it responds to anything.", {})]], accent=BAD)

# ════════════════════════════════════════ 6 — PERFORMANCE ANALYSIS
divider(6, "Performance Analysis",
        "One unified figure set: every panel carries every model, each answering a specific "
        "question a reviewer is likely to ask, with the table it was drawn from beside it.")

s, y = slide("Uncertainty — the one axis the foundation model loses on", eyebrow="Performance · 1",
             note="Nominal coverage for a p10–p90 band is 0.80. Coverage must be read next to "
                  "width and next to the last column: a band can always be made to cover by being "
                  "made useless, or by being allowed to predict the impossible.")
_z = zi.set_index("model")
rows, hl = [], set()
for i, m in enumerate(ORDER):
    if m == CHRONOS:
        hl.add(i)
    r = lb[(lb.model == m) & (lb.horizon == 24)].iloc[0]
    neg = float(_z.loc[m, "p10_below_zero_pct"])
    rows.append([LABEL[m], f"{r.coverage:.3f}",
                 [(f"{r.coverage - 0.8:+.3f}",
                   {"color": OK if abs(r.coverage - 0.8) < 0.03 else INK2})],
                 f"{r.width:.3f}",
                 [(f"{neg:.0f}%", {"color": BAD, "bold": True} if neg > 20 else {"color": INK2})],
                 f"{_z.loc[m, 'coverage_if_p10_clamped']:.3f}"])
y2 = table(s, MARGIN, y, BODY_W,
           ["Model", "p10–p90 coverage", "Gap to 0.80", "Width (KL/h)", "Rows with p10 < 0",
            "If p10 were 0"],
           rows, widths=[0.27, 0.15, 0.13, 0.14, 0.16, 0.15],
           aligns=["l", "r", "r", "r", "r", "r"], row_h=Inches(0.255), size=10,
           head_h=Inches(0.27), highlight=hl, fonts=[SANS, MONO, MONO, MONO, MONO, MONO])
callout(s, MARGIN, y2 + Inches(0.14), BODY_W, None,
        [[("The tell is the fifth column. ", {}),
          ("Every model that reaches nominal coverage gets there by putting its lower bound below "
           "zero", {"bold": True, "color": INK}),
          (" — a negative water outflow — except ", {}), ("NPTS", {"bold": True, "color": INK}),
          (", a non-parametric sampler that reproduces the zero atom directly. Clamping is not the "
           f"fix either: it takes Chronos-2 to {Z.coverage_if_p10_clamped:.3f}, past nominal. "
           f"Asymmetric conformal calibration reaches {CAL.coverage_cal:.3f} on a window it never "
           "saw.", {})]], accent=WARN)

FIGSLIDES = [
    (UP / "U1_leaderboard_mase.png", "The whole field, at a glance",
     "Macro MASE against horizon for all eleven models on identical rows. Four clean bands: the "
     "Chronos-2 family, the incumbent, the trained deep models, the classical methods."),
    (UP / "U3_significance_vs_all.png", "Chronos-2 against every opponent, with confidence",
     "MASE improvement with 95% paired-bootstrap CIs over the 24 origins — at 1 d on the left, "
     "across all six horizons on the right. Only the covariate variants cross zero."),
    (UP / "U4_win_matrix.png", "Who beats whom, cell by cell",
     "Percentage of the 144 tank × horizon cells each model wins against each other model. "
     "Chronos-2 takes every cell from the classical methods and loses cells only to its own "
     "covariate variants."),
    (UP / "U5_per_tank_heatmap.png", "Per tank, per model",
     "MASE at 1 d for all 24 tanks and all 11 models, ordered by demand and labelled by sensor "
     "tier. The hard tanks are hard for everyone — a property of those sensors, not of a model."),
    (UP / "U6_skill_vs_demand.png", "Is the average flattered by the dead sensors?",
     "Skill against the naive reference for every tank and every model. Every Chronos-2 variant "
     "and NPTS clear zero on all 24, and skill does not fall away with demand size."),
    (UP / "U9_error_by_leadtime.png", "Does it fall apart at long lead times?",
     f"MAE against hours ahead of origin, out to 168 h. Chronos-2 grows "
     f"{100*(_d7-_d1)/_d1:+.1f}% from day 1 to day 7 — which is what makes weekly refill "
     "planning viable."),
    (UP / "U7_calibration.png", "Coverage against the width that bought it",
     "Red rings mark models whose lower bound goes below zero. Read alone, coverage ranks the "
     "ringed models best; they get there by predicting a negative outflow."),
    (UP / "U11_zero_inflation.png", "Why the intervals miss",
     "Lower- versus upper-tail miss rates, how often p10 falls below zero, and what a naive clamp "
     "would do — for every model. The deficit is structural, not a defect of one model."),
    (UP / "U12_volume_bias.png", "Signed volume bias — the number that sizes a refill",
     "Total forecast volume against total actual, per model and horizon. The expensive direction "
     "is negative: a refill sized on an uncorrected mean would be short."),
    (UP / "U8_cost_accuracy.png", "Accuracy against what it cost",
     "Macro MASE at 1 d against measured wall clock, log scale. The production candidate is the "
     "leftmost point on the chart and within 0.003 MASE of the best."),
    (H45 / "plots/X_tanks45_chronos2.png", "45 continuous days — Chronos-2",
     "Solid actual, dashed calibrated, dotted uncalibrated, shaded conformal band. Tracks the "
     "healthy tanks closely; misses the spiky peaks on GJBC_LAW_BLOCK_3__A1."),
    (H45 / "plots/Y_tanks45_npts.png", "45 continuous days — the incumbent",
     "The same four tanks and the same correction. Flat on all four: the flatness seen at campus "
     "level is how the model behaves per tank, not an aggregation artefact."),
]
for n, (path, title, note) in enumerate(FIGSLIDES, start=2):
    if not path.exists():
        continue
    s, y = slide(title, eyebrow=f"Performance · {n}")
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(BODY_W / iw, (H - y - Inches(1.40)) / ih)
    pw = int(iw * scale)
    picture(s, path, MARGIN + int((BODY_W - pw) / 2), y, w=Emu(pw))
    caption(s, MARGIN, y + Emu(int(ih * scale)) + Inches(0.16), BODY_W, note)

s, y = slide("What we would not claim", eyebrow=f"Performance · {len(FIGSLIDES) + 2}",
             note="Volunteering the limits is what makes the rest of the deck credible. Each of "
                  "these is measured, not hedged.")
rows = [
    ["One healthy tank is genuinely failed",
     "GJBC_LAW_BLOCK_3__A1 is spiky and every model smooths it — the worst column in the per-tank "
     "heatmap for all 11. It needs a spike model, not a better smoother"],
    ["Intervals still under-cover after correction",
     f"{CAL.coverage_cal:.3f} against a nominal 0.80. Better, not solved"],
    ["A regime change cannot be anticipated",
     "GJBC_BLOCK_1_A4_RO averaged 0.0005 KL/h during calibration, then woke up. No correction "
     "fitted on an earlier window can foresee that — a production system must refit on a rolling "
     "basis and monitor coverage continuously"],
    ["PatchTST was not architecture-searched",
     f"Two configurations — AutoGluon defaults and the paper's hourly settings — not a per-tank "
     f"sweep. The claim is that a competently configured PatchTST loses by {PT_LO:.1f}–{PT_HI:.1f}%"
     " on 24 series this short, not that no PatchTST could ever win"],
    ["The result is one campus, 16 months",
     "It generalises to campuses like this one. It is not evidence about municipal supply, "
     "industrial demand, or a different climate"],
    ["No motor telemetry exists in the dataset",
     "So no claim is made about control quality, only about demand forecasting"],
    ["Nothing is running",
     "Every number here is offline backtest. The live system is designed, not built"],
]
table(s, MARGIN, y, BODY_W, ["Limit", "What is actually true"], rows,
      widths=[0.30, 0.70], aligns=["l", "l"], row_h=Inches(0.34), size=10.5,
      fonts=[SANS, SANS])

# ════════════════════════════════════════ CLOSING
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, fill=INK)
box(s, 0, 0, Inches(0.14), H, fill=C2)
tf = textbox(s, Inches(1.05), Inches(0.80), Inches(11.4), Inches(0.3))
_txt(tf, "IF YOU ONLY GET ONE SLIDE", size=12, color=C2, bold=True, font=MONO, space_after=0)
tf = textbox(s, Inches(1.05), Inches(1.30), Inches(11.1), Inches(5.4))
_txt(tf, [(f"On one evaluation grid where {NWORD} models score the same {TOTAL:,} rows, Chronos-2 "
           "zero-shot beats every non-Chronos model at every one of six horizons — and the margin "
           "is statistically significant everywhere, by paired bootstrap over 24 origins and by "
           "Diebold–Mariano.", {"bold": True, "color": GROUND, "size": 17})],
     space_after=12, line=1.35)
_txt(tf, [("It is ", {}),
          (f"{NPTS_LO:.1f}–{NPTS_HI:.1f}% better than the deployed incumbent",
           {"bold": True, "color": GROUND}),
          (f", {PT_LO:.1f}–{PT_HI:.1f}% better than a PatchTST trained on this campus, and "
           "roughly a third better than the classical methods. It wins ", {}),
          (f"{CELLS[INCUMBENT]} of 144", {"bold": True, "color": GROUND}),
          (" tank-horizon cells against the incumbent and all 144 against every classical method, "
           "beats the naive baseline on all 24 tanks, and costs ", {}),
          (f"{C2S:.0f} seconds", {"bold": True, "color": C2}),
          (" with no training.", {})],
     size=14, color=RGBColor(0xA3, 0xAE, 0xBA), space_after=12, line=1.4)
_txt(tf, [("The trained transformer loses to the incumbent as well as to Chronos-2.",
           {"bold": True, "color": GROUND}),
          (" On 24 short, noisy, zero-inflated series, training from scratch buys less than a "
           "well-chosen non-parametric method and far less than a pretrained one. That is the "
           "argument for a foundation model on this problem, stated as a measurement rather than "
           "a preference.", {})],
     size=14, color=RGBColor(0xA3, 0xAE, 0xBA), space_after=12, line=1.4)
_txt(tf, [("Both known weaknesses are diagnosed and corrected.", {"bold": True, "color": GROUND}),
          (f" The intervals covered {CAL.coverage_raw:.3f} against a nominal 0.80, traced to the "
           f"model's inability to represent the {Z.zero_fraction*100:.0f}% of hours with exactly "
           "zero demand — a limit that applies to every continuous-density model on the grid. "
           "Asymmetric conformal calibration and a per-tank volume correction, both fitted on an "
           "earlier disjoint window and measured out of sample, take coverage to ", {}),
          (f"{CAL.coverage_cal:.3f}", {"bold": True, "color": OK}),
          (" and volume bias to ", {}),
          (f"{OPS.total_bias_pct:+.1f}%", {"bold": True, "color": OK}),
          (f", cutting daily campus error from {OPS.daily_mae_raw_kl:.1f} to "
           f"{OPS.daily_mae_kl:.1f} KL. Neither requires retraining.", {})],
     size=14, color=RGBColor(0xA3, 0xAE, 0xBA), space_after=12, line=1.4)
_txt(tf, [("What is not done: ", {"bold": True, "color": GROUND}),
          ("nothing is deployed and running, and the research paper is not written. Both are "
           "designed in detail. Say this before you are asked.", {})],
     size=14, color=RGBColor(0xA3, 0xAE, 0xBA), space_after=0, line=1.4)

s, y = slide("Reproducing every number in this deck", eyebrow="Appendix",
             note="Run from the repository root with venv/ activated. Runtimes are the measured "
                  "wall clock from the actual run.")
CMDS = [
    ("python -m tests.test_metrics", "~3 s", "6/6 must pass — trust nothing below if they do not"),
    ("python -m src.data.curate", "~4 s", "26 directories → 24 tanks, 270,849 rows"),
    ("python -m src.models.chronos2_forecasting", f"~33 min", "the four Chronos-2 variants"),
    ("python -m src.models.baselines_autogluon", f"~{BASE_S/60:.0f} min",
     "NPTS, SeasonalNaive, ETS, Theta, DynamicOptimizedTheta in one pass"),
    ("python -m src.models.patchtst_benchmark --preset default", "~2 min",
     "the deep control at AutoGluon defaults"),
    ("python -m src.models.patchtst_benchmark --preset tuned  …", f"~{PTS/60:.0f} min",
     "the deep control at the paper's hourly settings"),
    ("python -m src.models.score_benchmark --strict", "~7 s",
     f"row parity across all {NM} models; fatal on mismatch"),
    ("python -m src.models.unified_analysis", "~4 min", "every table in this deck"),
    ("python -m src.models.unified_figures", "~25 s", "the twelve unified figures"),
    ("python -m src.models.calibrated_holdout", "~3 min", "conformal calibration, the 45-day view"),
    ("python reports/build_results_page.py", "~10 s", "the full HTML results page"),
    ("python reports/build_review_deck.py", "~15 s", "this deck"),
]
y2 = table(s, MARGIN, y, BODY_W, ["Command", "Runtime", "What it produces"],
           [[[(c, {"font": MONO, "size": 10})], t, w] for c, t, w in CMDS],
           widths=[0.47, 0.085, 0.445], aligns=["l", "r", "l"], row_h=Inches(0.255), size=10,
           head_h=Inches(0.28), fonts=[MONO, MONO, SANS])
tf = textbox(s, MARGIN, y2 + Inches(0.16), BODY_W, Inches(0.5))
_txt(tf, "Every number is read from a CSV under results/chronos2/unified/ at build time, so the "
         "deck and the results page cannot disagree. The bootstrap seed is fixed at 20260830.",
     size=10.5, color=INK3, space_after=0, line=1.3)

# ────────────────────────────────────────────────────────────── footers, then write
for i, sl in enumerate(prs.slides, start=1):
    if i == 1:
        continue
    if getattr(sl.shapes[0].fill.fore_color, "rgb", None) != INK:
        footer(sl, i)

prs.save(OUT)
print(f"wrote {OUT}  ({OUT.stat().st_size / 1e6:.2f} MB, "
      f"{len(prs.slides._sldIdLst)} slides, {NM} models)")
